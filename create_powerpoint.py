#!/usr/bin/env python3
"""
Create PowerPoint presentation with all 9 visualizations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

print("=" * 80)
print("📊 CREATING POWERPOINT PRESENTATION")
print("=" * 80)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define colors
BLUE = RGBColor(0, 81, 186)  # University Blue
RED = RGBColor(196, 30, 58)  # University Red
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(242, 242, 242)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(14), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    return slide

def add_content_slide(prs, title, subtitle, viz_file, talking_points, key_numbers):
    """Add content slide with visualization reference"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Header with title
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    header_frame = header.text_frame
    header_frame.text = title
    header_frame.paragraphs[0].font.size = Pt(40)
    header_frame.paragraphs[0].font.bold = True
    header_frame.paragraphs[0].font.color.rgb = BLUE

    # Subtitle
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(15), Inches(0.5))
        sub_frame = sub.text_frame
        sub_frame.text = subtitle
        sub_frame.paragraphs[0].font.size = Pt(20)
        sub_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    # Visualization placeholder with instructions
    viz_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(2), Inches(10), Inches(6)
    )
    viz_box.fill.solid()
    viz_box.fill.fore_color.rgb = LIGHT_GRAY
    viz_box.line.color.rgb = BLUE
    viz_box.line.width = Pt(3)

    text_frame = viz_box.text_frame
    text_frame.text = f"📊 VISUALIZATION\n\nOpen: {viz_file}\n\n(Interactive HTML chart)\n\nScreenshot and paste here,\nor demo live during presentation"
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = DARK_GRAY

    # Key numbers box
    if key_numbers:
        numbers_box = slide.shapes.add_textbox(Inches(11), Inches(2), Inches(4.5), Inches(6))
        numbers_frame = numbers_box.text_frame
        numbers_frame.word_wrap = True

        p = numbers_frame.paragraphs[0]
        p.text = "KEY NUMBERS"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RED

        for number in key_numbers:
            p = numbers_frame.add_paragraph()
            p.text = f"• {number}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(8)

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = f"TALKING POINTS:\n\n{talking_points}"

    return slide

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
print("\n[1/11] Creating title slide...")
add_title_slide(
    prs,
    "Strategic Revenue Optimization Plan",
    "Midwest State University Athletics | KAGR Case Competition 2025"
)

# ============================================================================
# SLIDE 2: Agenda
# ============================================================================
print("[2/11] Creating agenda slide...")
slide = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Today's Story: From Challenge to Solution"
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = BLUE

content = """
ACT 1: THE PROBLEM
    • The Challenge: $20.5M NCAA Settlement
    • Current State: Where we are today
    • The Gaps: Where we're behind

ACT 2: THE OPPORTUNITY
    • Women's Basketball: The paradox
    • 7 Strategic Initiatives

ACT 3: THE SOLUTION
    • Revenue Roadmap: $94M → $120M
    • Implementation Timeline
    • Return on Investment

FINALE: Executive Summary
"""

content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(13), Inches(6))
content_frame = content_box.text_frame
content_frame.text = content
content_frame.paragraphs[0].font.size = Pt(20)
content_frame.paragraphs[0].line_spacing = 1.5

# ============================================================================
# SLIDE 3: The Challenge
# ============================================================================
print("[3/11] Creating Challenge slide...")
add_content_slide(
    prs,
    "The Challenge",
    "NCAA vs. House Settlement Financial Impact",
    "docs/viz_01_challenge_gauge.html",
    """The NCAA vs. House settlement requires us to generate $20.5 million in new annual revenue.

This isn't optional—it's critical to maintaining our competitive excellence and supporting our student-athletes.

This gauge visualization shows the magnitude of the challenge we face. We need to be strategic, data-driven, and fan-focused in our approach.""",
    [
        "Required: $20.5M annually",
        "Timeline: Immediate",
        "Risk: High if not addressed",
        "Impact: All 6 sports"
    ]
)

# ============================================================================
# SLIDE 4: Current State
# ============================================================================
print("[4/11] Creating Current State slide...")
add_content_slide(
    prs,
    "Current State",
    "Midwest State Athletics Revenue Overview",
    "docs/viz_02_current_state_dashboard.html",
    """We currently generate $94.4 million annually across all athletic programs.

Ticket sales represent 42% of revenue at $39.8M, followed by concessions at 31% ($28.8M), parking at 15% ($13.9M), and merchandise at 12% ($11.8M).

We have 312 events across 6 sports with average attendance of 13,500 and 68.5% capacity utilization. There's room for growth.""",
    [
        "Total Revenue: $94.4M",
        "Ticket Sales: 42%",
        "Concessions: 31%",
        "Capacity: 68.5%",
        "Events: 312/year"
    ]
)

# ============================================================================
# SLIDE 5: The Gaps
# ============================================================================
print("[5/11] Creating Gap Analysis slide...")
add_content_slide(
    prs,
    "The Gaps",
    "Midwest State vs. Power 5 Conference Benchmarks",
    "docs/viz_03_gap_analysis.html",
    """When we compare ourselves to Power 5 conference peers like UT Austin, Ohio State, Michigan, Penn State, and Wisconsin, we see significant opportunities.

We're 38% below peers in corporate partnerships (9.2% vs 15%). We're 33% below in premium seating (8% vs 18%). Women's basketball capacity utilization is 21.5 points below average.

These gaps represent untapped revenue potential, not weaknesses. Our fans are here—we just need to engage them better.""",
    [
        "Corporate: 38% below peers",
        "Premium Seating: 56% below",
        "Women's BB: 33% below capacity",
        "Digital: 64% below leaders",
        "Source: NCAA Financial Database 2023-24"
    ]
)

# ============================================================================
# SLIDE 6: Women's Basketball Opportunity
# ============================================================================
print("[6/11] Creating Women's Basketball slide...")
add_content_slide(
    prs,
    "The Women's Basketball Paradox",
    "High Fan Interest, Low Attendance = Biggest Opportunity",
    "docs/viz_04_womens_bb_opportunity.html",
    """THIS is our biggest single opportunity.

Fan interest in women's basketball is 85 out of 100—nearly identical to men's basketball at 88. But our capacity utilization is only 43.5% compared to 84% for men.

Same arena. Same interest level. But 40 percentage points lower attendance.

If we simply bring women's basketball attendance to 60% capacity—still below men's—that's $4 million in annual revenue. Success stories at Iowa (174% growth) and South Carolina show this is achievable.""",
    [
        "Interest Score: 85/100",
        "Current Capacity: 43.5%",
        "Men's Capacity: 84%",
        "Opportunity: $4.0M/year",
        "Target: 60% capacity"
    ]
)

# ============================================================================
# SLIDE 7: 7 Strategic Initiatives
# ============================================================================
print("[7/11] Creating Initiatives slide...")
add_content_slide(
    prs,
    "7 Strategic Initiatives",
    "Balanced Portfolio: Quick Wins to Long-Term Investments",
    "docs/viz_05_initiative_bubbles.html",
    """We've identified 7 strategic initiatives, carefully balanced across timelines and effort levels.

QUICK WINS (3 months): Dynamic pricing and alumni programs - $5.1M combined, minimal investment.

SHORT-TERM (4-6 months): Merchandise expansion, women's basketball growth, digital platform - $8.7M combined.

STRATEGIC (9 months): Corporate partnerships - $7.5M, our largest single initiative.

LONG-TERM (12 months): Premium seating infrastructure - $2.8M annually, 3-year payback.

Together, they're diversified, achievable, and proven at peer institutions.""",
    [
        "Quick Wins: $5.1M (0-3mo)",
        "Short-term: $8.7M (3-6mo)",
        "Strategic: $7.5M (9mo)",
        "Long-term: $2.8M (12mo)",
        "Total: $24.1M",
        "Exceeds target by $3.6M"
    ]
)

# ============================================================================
# SLIDE 8: Revenue Waterfall
# ============================================================================
print("[8/11] Creating Revenue Waterfall slide...")
add_content_slide(
    prs,
    "Revenue Growth Roadmap",
    "From $94M to $120M: The Path Forward",
    "docs/viz_06_revenue_waterfall.html",
    """Here's our path from $94 million to $120 million—exceeding the $20.5M target by $4.6 million.

Each step builds on the last. Dynamic pricing adds $4.2M with minimal investment. Women's basketball growth adds $4.0M. Corporate partnerships add $7.5M—our largest initiative.

Combined with premium seating, digital transformation, merchandise expansion, and alumni programs, we reach $119.5M in total revenue.

That's $25.1M in new revenue, 122% of our target. We're not just meeting the challenge—we're exceeding it.""",
    [
        "Starting: $94.4M",
        "Target: $114.9M",
        "Projected: $119.5M",
        "New Revenue: $25.1M",
        "Exceeds by: $4.6M (22%)"
    ]
)

# ============================================================================
# SLIDE 9: Implementation Timeline
# ============================================================================
print("[9/11] Creating Timeline slide...")
add_content_slide(
    prs,
    "Implementation Roadmap",
    "18-Month Strategic Plan with Clear Milestones",
    "docs/viz_07_implementation_roadmap.html",
    """This is achievable in 18 months with clear milestones.

Q1 2025: Quick wins launch—dynamic pricing and alumni programs. $5.1M achieved.

Q2 2025: Phase 1 completes—add merchandise, women's basketball growth, digital platform. $14.6M total.

Q3 2025: Corporate partnerships ramp up. $17.6M total.

Q4 2025 & Beyond: Premium seating construction. Target exceeded at $20.5M+.

Each milestone is measurable, each initiative has an owner, and we have contingency plans.""",
    [
        "Q1: Quick wins ($5.1M)",
        "Q2: Phase 1 complete ($14.6M)",
        "Q3: Corporate ramp ($17.6M)",
        "Q4: Target achieved ($20.5M+)",
        "18 months: Full implementation"
    ]
)

# ============================================================================
# SLIDE 10: ROI Analysis
# ============================================================================
print("[10/11] Creating ROI slide...")
add_content_slide(
    prs,
    "Return on Investment",
    "Every Initiative is Above Break-Even",
    "docs/viz_08_roi_comparison.html",
    """Every single initiative delivers strong ROI.

Dynamic pricing has the highest ROI at 2,700%—requiring only $150K to generate $4.2M annually. 3-month payback.

Corporate partnerships require $500K but return $7.5M annually. 14-month payback, 1,400% ROI.

Even premium seating, our largest investment at $8.5M, delivers $2.8M annually with a 3-year payback and 280% 5-year ROI.

Average ROI across all initiatives: 723%. Every dollar invested returns over $7.""",
    [
        "Best ROI: 2,700% (Dynamic)",
        "Highest Return: $7.5M (Corporate)",
        "Average ROI: 723%",
        "Avg Payback: 10 months",
        "All above break-even"
    ]
)

# ============================================================================
# SLIDE 11: Executive Summary
# ============================================================================
print("[11/11] Creating Executive Summary slide...")
add_content_slide(
    prs,
    "Executive Summary",
    "The Complete Picture",
    "docs/viz_09_executive_summary.html",
    """In summary:

THE CHALLENGE: We face a $20.5M annual revenue gap from NCAA settlement obligations.

THE SOLUTION: 7 strategic initiatives spanning quick wins to long-term investments.

THE NUMBERS: We project $25.1M in new revenue—exceeding our target by 22%. Average ROI is 723%. Timeline is 18 months. Risk is low at 3.2 out of 10.

THE COMMITMENT: We maintain fan satisfaction above 4.0 throughout, enhance the gameday experience, and build sustainable competitive advantage.

We're ready to execute. We have the data, the plan, and the team to make this happen.""",
    [
        "Challenge: $20.5M",
        "Projected: $25.1M",
        "Exceeds by: 22%",
        "Initiatives: 7",
        "ROI: 723% average",
        "Timeline: 18 months",
        "Risk: LOW (3.2/10)"
    ]
)

# ============================================================================
# Save presentation
# ============================================================================
output_file = "KAGR_Presentation_Advanced_Visualizations.pptx"
prs.save(output_file)

print("\n" + "=" * 80)
print(f"✅ POWERPOINT CREATED: {output_file}")
print("=" * 80)
print("\n📊 11 slides created:")
print("   1. Title slide")
print("   2. Agenda")
print("   3. The Challenge")
print("   4. Current State")
print("   5. The Gaps")
print("   6. Women's Basketball Opportunity")
print("   7. 7 Strategic Initiatives")
print("   8. Revenue Waterfall")
print("   9. Implementation Timeline")
print("   10. ROI Analysis")
print("   11. Executive Summary")
print("\n💡 NEXT STEPS:")
print("   1. Open the PowerPoint file")
print("   2. Open each HTML visualization file")
print("   3. Screenshot the visualizations")
print("   4. Paste into the placeholder boxes")
print("   5. Adjust sizing as needed")
print("\n🎯 Your presentation is ready!")
print("=" * 80)
