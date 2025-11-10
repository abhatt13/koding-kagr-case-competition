#!/usr/bin/env python3
"""
KAGR Case Competition - PowerPoint Presentation Generator
Creates a comprehensive presentation based on the analysis and recommendations
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import Rectangle
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import io
from PIL import Image
import warnings
warnings.filterwarnings('ignore')

# Professional Color Palette
PRIMARY_BLUE = '#1f77b4'
PRIMARY_GREEN = '#2ca02c'
PRIMARY_RED = '#d62728'
PRIMARY_ORANGE = '#ff7f0e'
PRIMARY_PURPLE = '#9467bd'

EXCELLENT = '#2E7D32'
GOOD = '#66BB6A'
WARNING = '#FFA726'
CRITICAL = '#E53935'

# Configure matplotlib
plt.rcParams['font.family'] = 'sans-serif'
plt.rcParams['font.sans-serif'] = ['Arial', 'Helvetica', 'DejaVu Sans']
sns.set_palette("husl")

def create_chart_image(fig):
    """Convert matplotlib figure to image for PowerPoint"""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=300, bbox_inches='tight', facecolor='white')
    buf.seek(0)
    plt.close(fig)
    return buf

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Generating $20.5M in Incremental Revenue"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Strategic Initiatives for Midwest State Athletics"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Add author and date
    author_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
    author_frame = author_box.text_frame
    author_frame.text = "KODING with KAGR Case Competition\nNovember 2025"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(18)
    author_para.font.color.rgb = RGBColor(100, 100, 100)
    author_para.alignment = PP_ALIGN.CENTER

def add_challenge_slide(prs):
    """Slide 2: The Challenge & Context"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

    title = slide.shapes.title
    title.text = "The Challenge & Context"

    # Add content
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()

    # The Challenge section
    p = tf.paragraphs[0]
    p.text = "The Challenge:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(6)

    challenges = [
        "NCAA vs. House settlement: $20.5M annual impact (10 years)",
        "Revenue sharing starting 2025-26: up to 22% of revenue to athletes",
        "Bottom Line: Need 27% revenue growth to remain competitive"
    ]

    for challenge in challenges:
        p = tf.add_paragraph()
        p.text = challenge
        p.level = 1
        p.font.size = Pt(16)
        p.space_after = Pt(6)

    # Current Position section
    p = tf.add_paragraph()
    p.text = "\nCurrent Position:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(6)

    positions = [
        "Total Athletic Revenue: $94.4M (at NCAA Div I average)",
        "Strong programs: Football, Men's Basketball",
        "Question: Where do we find $20.5M+ without alienating fans?"
    ]

    for position in positions:
        p = tf.add_paragraph()
        p.text = position
        p.level = 1
        p.font.size = Pt(16)
        p.space_after = Pt(6)

def create_revenue_composition_chart():
    """Create revenue composition donut chart"""
    fig, ax = plt.subplots(figsize=(10, 8))

    revenue_sources = {
        'Ticket Sales': 39.83,
        'Concessions': 28.73,
        'Parking': 14.00,
        'Merchandise': 11.80
    }

    colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A']
    total = sum(revenue_sources.values())

    wedges, texts, autotexts = ax.pie(
        revenue_sources.values(),
        labels=revenue_sources.keys(),
        autopct=lambda pct: f'${pct/100*total:.1f}M\n({pct:.1f}%)',
        startangle=90,
        colors=colors,
        wedgeprops=dict(width=0.5, edgecolor='white', linewidth=3),
        textprops={'fontsize': 12, 'weight': 'bold'}
    )

    # Center circle
    centre_circle = plt.Circle((0,0), 0.70, fc='white')
    fig.gca().add_artist(centre_circle)

    ax.text(0, 0, f'Total Revenue\n${total:.1f}M',
            ha='center', va='center', fontsize=18, weight='bold',
            bbox=dict(boxstyle='round,pad=0.8', facecolor='lightgray', alpha=0.3))

    for autotext in autotexts:
        autotext.set_color('white')

    ax.set_title('Athletic Revenue Composition', fontsize=18, weight='bold', pad=20)
    plt.tight_layout()

    return fig

def create_womens_basketball_opportunity_chart():
    """Create the key Women's Basketball opportunity chart"""
    fig, ax1 = plt.subplots(figsize=(12, 8))

    sports = ['Football', "Men's BB", "Women's BB", 'Baseball', 'Softball', 'Volleyball']
    interest_scores = [95, 88, 85, 65, 58, 52]
    capacity_util = [86.7, 84.2, 43.5, 58.9, 61.8, 61.6]

    x = np.arange(len(sports))
    width = 0.35

    # Interest bars
    bars1 = ax1.bar(x - width/2, interest_scores, width, label='Fan Interest Score',
                    color='#1f77b4', edgecolor='black', linewidth=1.5, alpha=0.8)

    # Capacity utilization bars
    ax2 = ax1.twinx()
    colors_util = ['#2ca02c' if pct >= 70 else '#d62728' for pct in capacity_util]
    bars2 = ax2.bar(x + width/2, capacity_util, width,
                    label='Capacity Utilization %',
                    color=colors_util, edgecolor='black', linewidth=1.5, alpha=0.8)

    # Highlight Women's Basketball
    bars1[2].set_edgecolor('red')
    bars1[2].set_linewidth(4)
    bars2[2].set_edgecolor('red')
    bars2[2].set_linewidth(4)

    # Add annotation
    ax1.annotate('MAJOR OPPORTUNITY!\nHigh Interest (85)\nLow Attendance (43.5%)',
                 xy=(2, 85), xytext=(4.2, 90),
                 arrowprops=dict(arrowstyle='->', color='red', lw=3),
                 fontsize=14, color='red', weight='bold',
                 bbox=dict(boxstyle='round,pad=1', facecolor='yellow',
                          edgecolor='red', linewidth=3))

    ax1.set_xlabel('Sport', fontsize=14, weight='bold')
    ax1.set_ylabel('Fan Interest Score (0-100)', fontsize=13, weight='bold', color='#1f77b4')
    ax2.set_ylabel('Capacity Utilization %', fontsize=13, weight='bold', color='#2ca02c')
    ax1.set_title("Women's Basketball: High Interest, Low Attendance", fontsize=18, weight='bold', pad=20)
    ax1.set_xticks(x)
    ax1.set_xticklabels(sports, fontsize=11)

    ax1.legend(loc='upper left', fontsize=11)
    ax2.legend(loc='upper right', fontsize=11)
    ax1.grid(axis='y', alpha=0.3, linestyle='--')
    ax1.set_axisbelow(True)

    plt.tight_layout()
    return fig

def create_corporate_benchmark_chart():
    """Create corporate partnership benchmark comparison"""
    fig, ax = plt.subplots(figsize=(12, 8))

    schools = ['Wisconsin', 'Penn State', 'Michigan', 'Ohio State', 'UT Austin', 'Midwest State']
    corporate_pct = [14, 13, 14, 15, 16, 9.2]

    colors = ['#2ca02c' if x >= 13 else '#d62728' for x in corporate_pct]

    bars = ax.bar(schools, corporate_pct, color=colors, edgecolor='black',
                  linewidth=2, width=0.6, alpha=0.8)

    # Highlight Midwest State
    bars[-1].set_edgecolor('red')
    bars[-1].set_linewidth(3)
    bars[-1].set_hatch('//')

    # Average line
    avg = np.mean(corporate_pct[:-1])
    ax.axhline(avg, color='blue', linestyle='--', linewidth=2,
               label=f'Peer Average: {avg:.1f}%', alpha=0.7)

    # Add values on bars
    for i, (school, val) in enumerate(zip(schools, corporate_pct)):
        ax.text(i, val + 0.3, f'{val:.1f}%', ha='center', va='bottom',
                fontsize=12, weight='bold')

    # Gap annotation
    gap = avg - 9.2
    ax.annotate(f'{gap:.1f}% below\npeer average',
                xy=(5, 9.2), xytext=(4, 6),
                arrowprops=dict(arrowstyle='->', color='red', lw=2),
                fontsize=12, color='red', weight='bold',
                bbox=dict(boxstyle='round,pad=0.5', facecolor='yellow', alpha=0.7))

    ax.set_ylabel('Corporate Partnerships (% of total attendance)', fontsize=13, weight='bold')
    ax.set_title('Corporate Partnership Gap vs. Power 5 Peers', fontsize=18, weight='bold', pad=20)
    ax.set_ylim(0, max(corporate_pct) * 1.2)
    ax.legend(fontsize=11)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    ax.set_axisbelow(True)

    plt.tight_layout()
    return fig

def create_revenue_waterfall_chart():
    """Create revenue growth waterfall chart"""
    fig, ax = plt.subplots(figsize=(16, 9))

    initiatives = ['Current\nRevenue', 'Dynamic\nPricing', 'Women BB\nGrowth',
                   'Corporate\nPartnerships', 'Merchandise', 'Digital\nPlatform',
                   'Premium\nSeating', 'Off-Peak', 'Projected\nRevenue']

    values = [94.36, 3.8, 4.4, 7.5, 1.9, 2.8, 2.8, 0.9, 0]

    cumulative = [94.36]
    for i in range(1, len(values)-1):
        cumulative.append(cumulative[-1] + values[i])
    cumulative.append(cumulative[-1])

    colors = ['#1f77b4'] + ['#2ca02c'] * 7 + ['#ff7f0e']

    for i, (init, val, cum) in enumerate(zip(initiatives, values, cumulative)):
        if i == 0 or i == len(initiatives) - 1:
            ax.bar(i, cum, color=colors[i], edgecolor='black', linewidth=2, width=0.7)
            ax.text(i, cum + 2, f'${cum:.1f}M', ha='center', va='bottom',
                    fontsize=13, weight='bold')
        else:
            bottom = cumulative[i-1]
            ax.bar(i, val, bottom=bottom, color=colors[i],
                   edgecolor='black', linewidth=2, width=0.7)
            ax.text(i, bottom + val/2, f'+${val:.1f}M', ha='center', va='center',
                    fontsize=12, weight='bold', color='white',
                    bbox=dict(boxstyle='round,pad=0.3', facecolor='black', alpha=0.6))

            if i < len(initiatives) - 1:
                ax.plot([i+0.35, i+0.65], [cumulative[i], cumulative[i]],
                        'k--', linewidth=1.5, alpha=0.5)

    ax.set_xticks(range(len(initiatives)))
    ax.set_xticklabels(initiatives, fontsize=11, weight='bold')
    ax.set_ylabel('Revenue (Millions $)', fontsize=14, weight='bold')
    ax.set_title('Revenue Growth Roadmap: $94.4M → $118.5M (+$24.1M, 24% above target)',
                 fontsize=16, weight='bold', pad=20)

    # Target line
    target = 94.36 + 20.5
    ax.axhline(target, color='#d62728', linestyle='--', linewidth=2.5,
               label=f'NCAA Settlement Target: ${target:.1f}M', alpha=0.7)

    ax.legend(fontsize=12, loc='upper left')
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    ax.set_axisbelow(True)
    ax.set_ylim(0, cumulative[-1] * 1.15)

    plt.tight_layout()
    return fig

def add_slide_with_chart(prs, title_text, chart_func, bullet_points=None):
    """Add a slide with a chart and optional bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # Create and add chart
    fig = chart_func()
    img_stream = create_chart_image(fig)

    if bullet_points:
        # Add chart on the right
        slide.shapes.add_picture(img_stream, Inches(4.5), Inches(1.2), width=Inches(5))

        # Add bullet points on the left
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3.5), Inches(5))
        tf = text_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(14)
            p.space_after = Pt(8)
    else:
        # Add chart centered and larger
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.2), width=Inches(9))

def add_initiatives_summary_slide(prs):
    """Slide: Full Initiative Portfolio"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Seven Strategic Initiatives"

    # Add table
    rows = 9
    cols = 5
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(0.5)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(2)
    table.columns[4].width = Inches(1.5)

    # Header row
    headers = ['#', 'Initiative', 'Annual Impact', 'Implementation Cost', 'Net Impact']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    data = [
        ['1', 'Dynamic Pricing', '$4.8M', '$1.0M', '$3.8M'],
        ['2', "Women's BB Growth", '$4.6M', '$200K', '$4.4M'],
        ['3', 'Corporate Partnerships', '$8.2M', '$700K', '$7.5M'],
        ['4', 'Premium Seating Expansion', '$4.0M', '$1.2M*', '$2.8M'],
        ['5', 'Digital Platform', '$3.2M', '$400K', '$2.8M'],
        ['6', 'Merchandise Optimization', '$2.1M', '$200K', '$1.9M'],
        ['7', 'Alumni Engagement', '$1.0M', '$100K', '$0.9M'],
        ['', 'TOTAL', '$27.9M', '$3.8M', '$24.1M'],
    ]

    for i, row_data in enumerate(data):
        for j, value in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            if i == len(data) - 1:  # Total row
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 220, 220)

def add_implementation_timeline_slide(prs):
    """Slide: Implementation Flexibility"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Implementation Flexibility - Phased Approach"

    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()

    # Tier 1
    p = tf.paragraphs[0]
    p.text = "Tier 1 - Start Here (Months 0-3): $15.7M"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 51)

    tier1_items = [
        "Dynamic Pricing ($3.8M net)",
        "Women's Basketball Growth ($4.4M net)",
        "Corporate Partnership Expansion ($7.5M net)"
    ]

    for item in tier1_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(14)
        p.space_after = Pt(6)

    # Tier 2
    p = tf.add_paragraph()
    p.text = "\nTier 2 - Quick Wins (Months 3-9): $20.4M Total"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 102, 0)

    tier2_items = [
        "Digital Platform ($2.8M net)",
        "Merchandise Optimization ($1.9M net)"
    ]

    for item in tier2_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(14)
        p.space_after = Pt(6)

    # Tier 3
    p = tf.add_paragraph()
    p.text = "\nTier 3 - Long-term (Months 9-18): $24.1M Total"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    tier3_items = [
        "Premium Seating ($2.8M net, capital intensive)",
        "Alumni Engagement ($0.9M net)"
    ]

    for item in tier3_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(14)
        p.space_after = Pt(6)

def add_conclusion_slide(prs):
    """Final slide: Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "The Bottom Line"

    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Summary:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    summary_items = [
        "Challenge: Generate $20.5M to offset House settlement",
        "Solution: 7 data-driven initiatives",
        "Result: $24.1M net revenue (24% above target)",
        "Timeline: 12-18 months to full implementation"
    ]

    for item in summary_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(16)
        p.space_after = Pt(10)

    # Next Steps
    p = tf.add_paragraph()
    p.text = "\nImmediate Next Steps:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_before = Pt(15)

    next_steps = [
        "Week 1-2: Approve Tier 1 initiatives, RFP for dynamic pricing vendor",
        "Months 1-3: Launch women's basketball campaign, pilot dynamic pricing",
        "Months 6-18: Full deployment, measure and optimize"
    ]

    for step in next_steps:
        p = tf.add_paragraph()
        p.text = step
        p.level = 1
        p.font.size = Pt(16)
        p.space_after = Pt(10)

def main():
    """Main function to create the presentation"""
    print("Creating PowerPoint presentation...")

    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add slides
    print("Adding Slide 1: Title Slide")
    add_title_slide(prs)

    print("Adding Slide 2: Challenge & Context")
    add_challenge_slide(prs)

    print("Adding Slide 3: Revenue Composition")
    add_slide_with_chart(prs, "Current State: Revenue Composition",
                         create_revenue_composition_chart)

    print("Adding Slide 4: Women's Basketball Opportunity")
    womens_bb_bullets = [
        "Survey shows 85 interest score (vs. Men's 88)",
        "Only 43.5% capacity utilization",
        "Same 12,000-seat arena as Men's",
        "Gap = 3,988 empty seats/game",
        "Opportunity: $4.4M annually"
    ]
    add_slide_with_chart(prs, "Critical Finding #1: Women's Basketball",
                         create_womens_basketball_opportunity_chart,
                         womens_bb_bullets)

    print("Adding Slide 5: Corporate Partnership Gap")
    corporate_bullets = [
        "Midwest State: 9.2% corporate",
        "Peer average: 15%",
        "38% below benchmark",
        "Opportunity: +100K corporate attendees",
        "Projected impact: $7.5M annually"
    ]
    add_slide_with_chart(prs, "Critical Finding #2: Corporate Partnership Gap",
                         create_corporate_benchmark_chart,
                         corporate_bullets)

    print("Adding Slide 6: Revenue Waterfall")
    add_slide_with_chart(prs, "Full Initiative Portfolio",
                         create_revenue_waterfall_chart)

    print("Adding Slide 7: Initiatives Summary Table")
    add_initiatives_summary_slide(prs)

    print("Adding Slide 8: Implementation Timeline")
    add_implementation_timeline_slide(prs)

    print("Adding Slide 9: Conclusion")
    add_conclusion_slide(prs)

    # Save presentation
    output_file = '/home/user/koding-kagr-case-competition/KAGR_Case_Competition_Presentation.pptx'
    prs.save(output_file)
    print(f"\n✅ Presentation saved successfully to: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
